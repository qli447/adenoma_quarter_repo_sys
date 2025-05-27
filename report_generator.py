import mysql.connector
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from datetime import date, datetime
import streamlit as st
import os
import re

def generate_report(facility_name: str) -> str:
    # ====== 1. Determine previous quarter date range ======
    today = date.today()
    year, month = today.year, today.month
    if month <= 3:
        prev_quarter = 4
        prev_year = year - 1
    else:
        prev_quarter = (month - 1) // 3
        prev_year = year

    start_month = (prev_quarter - 1) * 3 + 1
    end_month = start_month + 2
    quarter_months = [f"{prev_year}-{str(m).zfill(2)}" for m in range(start_month, end_month + 1)]
    start_date = f"{prev_year}-{str(start_month).zfill(2)}-01"
    end_date = f"{prev_year}-{str(end_month).zfill(2)}-31"
    quarter_tag = f"Q{prev_quarter}_{prev_year}"

    # ====== 2. Connect to DB ======
    conn = mysql.connector.connect(
            host=st.secrets["DB_HOST"],
            user=st.secrets["DB_USER"],
            password=st.secrets["DB_PASSWORD"],
            database=st.secrets["DB_NAME"],
            port=3306
        )

    # ====== 3. Query Data ======
    query = f"""
    WITH tmp AS (
        SELECT DISTINCT
            a.gopath_id,
            SUBSTR(a.signout_date, 1, 7) AS signout,
            b.category AS diag,
            a.attending_physician,
            b.gender AS gender,
            a.facility_name
        FROM
            rpt_ap_case a
        JOIN
            rpt_ap_diag_general b ON a.gopath_id = b.gopath_id
        WHERE
            a.signout_date BETWEEN '{start_date}' AND '{end_date}'
            AND a.facility_name LIKE "%{facility_name}%"
    ),
    summary AS (
        SELECT
            signout,
            diag,
            attending_physician,
            gender,
            facility_name,
            COUNT(*) AS ct
        FROM
            tmp
        GROUP BY
            signout, diag, attending_physician, gender, facility_name
    )
    SELECT
        attending_physician,
        diag,
        gender,
        facility_name,
    """ + ",\n".join([f"    SUM(CASE WHEN signout = '{m}' THEN ct ELSE 0 END) AS `{m}`" for m in quarter_months]) + """
    FROM
        summary
    GROUP BY
        attending_physician, diag, gender, facility_name
    ORDER BY
        attending_physician ASC, gender DESC, diag ASC;
    """

    df = pd.read_sql(query, conn)
    conn.close()

    if df.empty:
        return None

    actual_facility_name = df["facility_name"].iloc[0]
    group_columns = ["attending_physician", "diag", "gender", "facility_name"]
    existing_months = [col for col in quarter_months if col in df.columns]
    df = df[group_columns + existing_months]

    # ====== 4. Load PPT Template ======
    template_path = os.path.join(os.path.dirname(__file__), "Gopath template.pptx")
    prs = Presentation(template_path)
    blank_slide_layout = prs.slide_layouts[6]

    # ====== 5. Update intro slide ======
    title_text = f"{actual_facility_name} - Case Review\n(Breakdown by Doctors)"
    subtitle_text = f"Received Cases: {quarter_tag}, GPIS LIS"
    report_date = datetime.today().strftime("%B %d, %Y")
    subtitle_full = f"{subtitle_text}\nDate Reported: {report_date}"
    title_set, subtitle_set = False, False

    intro_slide = prs.slides[0]
    for shape in intro_slide.shapes:
        if shape.has_text_frame:
            shape.text_frame.clear()
            p = shape.text_frame.paragraphs[0]
            run = p.add_run()
            if not title_set:
                run.text = title_text
                run.font.bold = True
                run.font.size = Pt(30)
                run.font.name = "Arial"
                title_set = True
            elif not subtitle_set:
                run.text = subtitle_full
                run.font.size = Pt(20)
                run.font.name = "Arial"
                subtitle_set = True

    # ====== 6. Slide layout helpers ======
    footer_text = "contact@gopathlabs.com ∙ 855-GOPATH9  www.gopathdx.com   1000 Corporate Grove Dr, Buffalo Grove, IL 60089"
    slide_width_inches = prs.slide_width.inches
    slide_height_inches = prs.slide_height.inches

    def create_slide(title_text, doctor, display_df):
        slide = prs.slides.add_slide(blank_slide_layout)
        prs.slides._sldIdLst.insert(1, prs.slides._sldIdLst[-1])
        for shape in list(slide.shapes):
            if shape.is_placeholder:
                slide.shapes._spTree.remove(shape._element)

        title_box = slide.shapes.add_textbox(Inches(0.3), Inches(0.2), Inches(6), Inches(0.4))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title_text.title()
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 51, 51)

        right_box = slide.shapes.add_textbox(Inches(6.3), Inches(0.2), Inches(3), Inches(0.6))
        tf2 = right_box.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = f"Facility: {actual_facility_name}\nDoctor: {doctor}"
        p2.font.size = Pt(10)

        rows, cols = display_df.shape
        col_widths = [2.5 if col == "diag" else 0.5 for col in display_df.columns]
        total_width = sum(col_widths)
        table_left = Inches((slide_width_inches - total_width) / 2)
        table_top = Inches(0.9)
        row_height_inches = 0.3
        max_table_height = 5.5
        table_height = min(max_table_height, (rows + 1) * row_height_inches)

        table = slide.shapes.add_table(rows + 1, cols, table_left, table_top,
                                       Inches(total_width), Inches(table_height)).table

        for i, width in enumerate(col_widths):
            table.columns[i].width = Inches(width)

        for i, col_name in enumerate(display_df.columns):
            cell = table.cell(0, i)
            cell.text = col_name
            for para in cell.text_frame.paragraphs:
                for run in para.runs:
                    run.font.size = Pt(9)
                    run.font.bold = True

        for i in range(rows):
            for j in range(cols):
                cell = table.cell(i + 1, j)
                value = display_df.iat[i, j]
                cell.text = str(int(value)) if isinstance(value, (int, float)) else str(value)
                for para in cell.text_frame.paragraphs:
                    for run in para.runs:
                        run.font.size = Pt(8)

        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(6.9), Inches(9), Inches(0.3))
        tf = txBox.text_frame
        p = tf.add_paragraph()
        p.text = footer_text
        p.font.size = Pt(8)

    # ====== 7. Generate slides per doctor ======
    for doctor, doc_df in df.groupby("attending_physician"):
        overall_df = doc_df.drop(columns=["gender", "facility_name"])
        overall_df = overall_df.groupby(["attending_physician", "diag"], as_index=False).sum()
        display_df = overall_df.drop(columns=["attending_physician"])
        create_slide("Diagnosis results, overall", doctor, display_df)

        for gender, gender_df in doc_df.groupby("gender"):
            display_df = gender_df.drop(columns=["attending_physician", "gender", "facility_name"])
            gender_label = gender if gender and str(gender).strip() else "Other"
            create_slide(f"Diagnosis results, {gender_label} patients", doctor, display_df)

    # ====== 8. Save & return path ======
    safe_facility_name = re.sub(r'[^\w\s-]', '', actual_facility_name).strip().replace(' ', '_')
    filename = f"{safe_facility_name}-Report-{quarter_tag}.pptx"
    output_dir = os.path.join(os.path.dirname(__file__), "output")
    os.makedirs(output_dir, exist_ok=True)
    output_path = os.path.join(output_dir, filename)
    prs.save(output_path)

    return output_path
